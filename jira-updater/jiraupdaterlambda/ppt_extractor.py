import base64
import os
from io import BytesIO
from pptx import Presentation
from pdf2image import convert_from_path
import subprocess
import json
from PIL import Image
import re

import subprocess
import os

def convert_pdf_to_images(pdf_path, output_dir):
    # Ensure the output directory exists.
    os.makedirs(output_dir, exist_ok=True)

    # Use the PDF file's base name (without extension) as the prefix.
    base_name = os.path.splitext(os.path.basename(pdf_path))[0]
    output_prefix = os.path.join(output_dir, base_name)

    # Convert PDF pages to PNG images using pdftoppm.
    command = ["pdftoppm", pdf_path, output_prefix, "-png"]
    subprocess.run(command, check=True)

    # Get PDF info using pdfinfo to extract the number of pages.
    info_command = ["pdfinfo", pdf_path]
    result = subprocess.run(info_command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)

    # Parse the output to find the line that starts with "Pages:" and extract the page count.
    page_count = None
    for line in result.stdout.splitlines():
        if line.startswith("Pages:"):
            page_count = int(line.split(":")[1].strip())
            break

    if page_count is None:
        raise RuntimeError("Could not determine page count from pdfinfo output.")

    # List generated images in the output directory that match the base name.
    image_files = []
    for file in os.listdir(output_dir):
        if file.startswith(base_name) and file.endswith(".png"):
            image_files.append(file)

    # Sort the image files by page number extracted from their filenames.
    # Expected filename format: <base_name>-<number>.png
    def extract_page_number(filename):
        match = re.search(rf"{re.escape(base_name)}-(\d+)\.png", filename)
        if match:
            return int(match.group(1))
        return 0

    image_files.sort(key=extract_page_number)

    # Open images with PIL and store them in a list.
    pages = []
    for filename in image_files:
        image_path = os.path.join(output_dir, filename)
        image = Image.open(image_path)
        pages.append(image)

    if len(pages) != page_count:
        print(f"Warning: pdfinfo reported {page_count} pages, but found {len(pages)} images.")

    return pages

def decode_base64_to_pptx(base64_string, output_path):
    """
    Decodes a base64 string and saves it as a .pptx file.
    """
    with open(output_path, "wb") as file:
        file.write(base64.b64decode(base64_string))


def convert_pptx_to_pdf(pptx_path, pdf_path):
    """
    Converts a .pptx PowerPoint file to a .pdf file using LibreOffice in headless mode.
    
    This function works on Linux provided that LibreOffice is installed.
    It runs the command:
      libreoffice --headless --convert-to pdf <pptx_path> --outdir <output_dir>
    and then renames the generated PDF to the desired pdf_path.
    """
    try:
        output_dir = os.path.dirname(pdf_path)
        # Run LibreOffice in headless mode to convert the PPTX to PDF.
        subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf', pptx_path, '--outdir', output_dir],
            check=True
        )
        # LibreOffice names the PDF with the same base name as the PPTX.
        generated_pdf = os.path.join(output_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
        if generated_pdf != pdf_path:
            os.rename(generated_pdf, pdf_path)
        print(f"✅ PowerPoint successfull converted to PDF: {pdf_path}")
    except Exception as e:
        print(f"❌ Error converting PowerPoint to PDF: {e}")

def extract_pdf_pages_to_png(pdf_path, slide_numbers, output_folder):
    """
    Extracts specified slides from a PDF and saves them as PNG images.
    """
    pages = convert_from_path(pdf_path)
    for page_num in slide_numbers:
        if 1 <= page_num <= len(pages):
            image_path = os.path.join(output_folder, f"slide_{page_num}.png")
            pages[page_num - 1].save(image_path, "PNG")
            print(f"Saved slide {page_num} as {image_path}")

def parse_configuration_slide(slide):
    """
    Extracts configuration values from the first slide.
    Expected format in the slide (example):
        Functionalities: 7, 8, 9
        Scopes: 10, 11, 12
        VISA: https://example.com/visa_link
    Returns a dictionary mapping the keys to lists of slide numbers (for Functionalities/Scopes)
    or to a string (for VISA).
    """
    config_map = {}
    all_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
    for line in all_text.splitlines():
        parts = line.split(":")
        if len(parts) >= 2:
            key = parts[0].strip()
            # Join in case the value itself contains colons
            value = ":".join(parts[1:]).strip()
            if key in ["Functionalities", "Scopes"]:
                try:
                    numbers = [int(x.strip()) for x in value.split(",") if x.strip().isdigit()]
                    config_map[key] = numbers
                except ValueError:
                    pass
            elif key == "VISA":
                config_map["VISA"] = value
    return config_map

def extract_functionalities_data(slide):
    """
    Extracts table data from a functionalities slide.
    Expected columns:
      - 'Numéro'
      - 'ID Jira'
      - 'Nom'
      - 'Hypothèses de bénéfices'
      - 'Critères d’acceptance'
      - 'Sizing SI (PI)'
    Returns a list of dictionaries with the above keys.
    """
    expected_columns = ['Numéro', 'ID Jira', 'Nom', 'Hypothèses de bénéfices', 'Critères d’acceptance', 'Sizing SI (PI)']
    functionalities = []
    for shape in slide.shapes:
        if hasattr(shape, "has_table") and shape.has_table:
            table = shape.table
            rows = []
            for r in table.rows:
                row_data = [cell.text.strip() for cell in r.cells]
                rows.append(row_data)
            # Assume the table has a header row followed by data rows.
            if len(rows) >= 2:
                data_rows = rows[1:]
                for row in data_rows:
                    if len(row) == len(expected_columns):
                        row_dict = {col: value for col, value in zip(expected_columns, row)}
                        functionalities.append(row_dict)
    return functionalities

def extract_scopes_data(slide, slide_number):
    """
    Extracts table data from a scopes slide.
    Expected columns:
      - 'Numéro'
      - 'Scopes'
      - 'Applications'
      - 'Referents'
      - 'Impacts / Architecture'
    Returns a list of dictionaries with the above keys.
    Adds the key 'source_slide' to track the slide number.
    """
    expected_columns = ['Numéro', 'Scopes', 'Applications', 'Referents', 'Impacts / Architecture']
    scopes = []
    for shape in slide.shapes:
        if hasattr(shape, "has_table") and shape.has_table:
            table = shape.table
            rows = []
            for r in table.rows:
                row_data = [cell.text.strip() for cell in r.cells]
                rows.append(row_data)
            if len(rows) >= 2:
                data_rows = rows[1:]
                for row in data_rows:
                    if len(row) == len(expected_columns):
                        row_dict = {col: value for col, value in zip(expected_columns, row)}
                        # Still record the original slide number (optional)
                        row_dict["source_slide"] = slide_number
                        scopes.append(row_dict)
    return scopes

def parse_impacts_architecture(text):
    """
    Parses the text in the format:
        "Impacts: <numbers separated by comma>\nArchitectures: <numbers separated by comma>"
    and returns a tuple (list_of_impact_numbers, list_of_architecture_numbers).
    """
    impacts = []
    architectures = []
    for line in text.split("\n"):
        if line.startswith("Impacts:"):
            raw = line[len("Impacts:"):].strip()
            if raw:
                impacts = [int(x.strip()) for x in raw.split(",") if x.strip().isdigit()]
        elif line.startswith("Architectures:"):
            raw = line[len("Architectures:"):].strip()
            if raw:
                architectures = [int(x.strip()) for x in raw.split(",") if x.strip().isdigit()]
    return impacts, architectures

def associate_impacts_architecture(functionalities, scopes):
    """
    Creates a mapping for each functionality (based on 'Numéro') to its corresponding 
    'Impacts / Architecture' extracted from scopes slides.
    Instead of keeping the raw text, this function splits it into separate keys:
      - "Impacts:" as a list of slide numbers to convert to PNG.
      - "Architectures:" as a list of slide numbers to convert to PNG.
    Returns a list of dictionaries.
    """
    association = []
    for func in functionalities:
        num = func.get("Numéro")
        combined_impacts = []
        combined_architectures = []
        for scope in scopes:
            if scope.get("Numéro") == num:
                raw_value = scope.get("Impacts / Architecture", "")
                impacts, architectures = parse_impacts_architecture(raw_value)
                combined_impacts.extend(impacts)
                combined_architectures.extend(architectures)
        # Remove duplicates and sort the numbers
        association.append({
            "Numéro": num,
            "Impacts:": sorted(list(set(combined_impacts))),
            "Architectures:": sorted(list(set(combined_architectures)))
        })
    return association

def process_pptx(base64_pptx):
    """
    Process a base64 PowerPoint:
      - Reads configuration from the first slide (slide numbers for functionalities and scopes, and the VISA link).
      - Extracts table data from functionalities and scopes slides with structured columns.
      - Creates an association table for each functionality's 'Numéro' with its corresponding 'Impacts / Architecture'.
      - Converts the PPTX to a PDF.
      - Uses the impacts_architecture_association list to extract PNG images from the PDF.
        Each PNG is saved with a file name of the feature including an indication if it's impact or architecture.
      - Structures all extracted information in a JSON object.
    """

    pptx_path = os.path.join("/tmp/", "presentation.pptx")
    pdf_path = os.path.join("/tmp/", "presentation.pdf")
    json_file = os.path.join("/tmp/", "extracted_data.json")

    decode_base64_to_pptx(base64_pptx, pptx_path)
    convert_pptx_to_pdf(pptx_path, pdf_path)

    prs = Presentation(pptx_path)
    config_slide = prs.slides[0]
    config_map = parse_configuration_slide(config_slide)

    extracted_data = {}
    extracted_data["config"] = config_map
    functionalities_data = []
    scopes_data = []

    # Extract functionalities tables
    if "Functionalities" in config_map:
        for slide_number in config_map["Functionalities"]:
            slide_index = slide_number - 1  # Convert to 0-index
            if slide_index < len(prs.slides):
                func_data = extract_functionalities_data(prs.slides[slide_index])
                if func_data:
                    functionalities_data.extend(func_data)
                else:
                    print(f"No table data found in functionalities slide {slide_number}.")
            else:
                print(f"Slide {slide_number} for functionalities not found in the presentation.")
    extracted_data["functionalities"] = functionalities_data


    # Extract scopes tables (adding source_slide)
    if "Scopes" in config_map:
        for slide_number in config_map["Scopes"]:
            slide_index = slide_number - 1
            if slide_index < len(prs.slides):
                scope_data = extract_scopes_data(prs.slides[slide_index], slide_number)
                if scope_data:
                    scopes_data.extend(scope_data)
                else:
                    print(f"No table data found in scopes slide {slide_number}.")
            else:
                print(f"Slide {slide_number} for scopes not found in the presentation.")
    extracted_data["scopes"] = scopes_data

    # Create association table: for each functionality 'Numéro', associate its impacts and architectures
    extracted_data["impacts_architecture_association"] = associate_impacts_architecture(functionalities_data, scopes_data)

    # ---- New Part: Extract PNG images using the association list ----
    try:
        #pages = convert_from_path(pdf_path=pdf_path, poppler_path="/usr/bin/pdftoppm")
        pages = convert_pdf_to_images(pdf_path, "/tmp/output")

        for assoc in extracted_data["impacts_architecture_association"]:
            feature_num = assoc.get("Numéro")
            # Process Impact images: the numbers here represent slide numbers to convert to PNG.
            for slide_num in assoc.get("Impacts:", []):
                if 1 <= slide_num <= len(pages):
                    image_path = os.path.join("/tmp/output", f"{feature_num}_impact_slide{slide_num}.png")
                    pages[slide_num - 1].save(image_path, "PNG")
                    print(f"Saved impact PNG: {image_path}")
            # Process Architecture images.
            for slide_num in assoc.get("Architectures:", []):
                if 1 <= slide_num <= len(pages):
                    image_path = os.path.join("/tmp/output", f"{feature_num}_architecture_slide{slide_num}.png")
                    pages[slide_num - 1].save(image_path, "PNG")
                    print(f"Saved architecture PNG: {image_path}")

    except Exception as e:
        print("Error:", e)

    return extracted_data
